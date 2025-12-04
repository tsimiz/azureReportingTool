"""Module for generating PowerPoint presentations from analysis results."""

import logging
from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class PowerPointGenerator:
    """Generates PowerPoint presentations from Azure analysis results."""

    def __init__(self):
        """Initialize PowerPoint generator."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        logger.info("PowerPoint generator initialized")

    def add_title_slide(self, title: str = "Azure Environment Analysis Report"):
        """Add title slide to presentation."""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = "Comprehensive analysis against Microsoft best practices"
        
        logger.info("Added title slide")

    def add_executive_summary_slide(self, summary: str):
        """Add executive summary slide."""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Split summary into paragraphs
        paragraphs = summary.split('\n\n')
        for i, para in enumerate(paragraphs):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = para.strip()
            p.font.size = Pt(12)
            p.level = 0
        
        logger.info("Added executive summary slide")

    def add_resource_overview_slide(self, resources: Dict[str, List[Dict[str, Any]]]):
        """Add resource overview slide with counts."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Azure Environment Overview"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add resource counts
        resource_counts = [
            f"Resource Groups: {len(resources.get('resource_groups', []))}",
            f"Virtual Machines: {len(resources.get('virtual_machines', []))}",
            f"Storage Accounts: {len(resources.get('storage_accounts', []))}",
            f"Network Security Groups: {len(resources.get('network_security_groups', []))}",
            f"Virtual Networks: {len(resources.get('virtual_networks', []))}",
            f"Total Resources: {len(resources.get('all_resources', []))}"
        ]
        
        for i, count in enumerate(resource_counts):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = count
            p.font.size = Pt(18)
            p.font.bold = True
            p.level = 0
        
        logger.info("Added resource overview slide")

    def add_findings_slide(self, resource_type: str, analysis: Dict[str, Any]):
        """Add slide for findings of a specific resource type."""
        if not analysis or 'findings' not in analysis or not analysis['findings']:
            return
        
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{resource_type.replace('_', ' ').title()} - Findings"
        
        # Add overall score if available
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        if 'overall_score' in analysis:
            score = analysis['overall_score']
            p = text_frame.paragraphs[0]
            p.text = f"Overall Score: {score}/10"
            p.font.size = Pt(16)
            p.font.bold = True
            
            # Color code the score
            if score >= 8:
                p.font.color.rgb = RGBColor(0, 128, 0)  # Green
            elif score >= 6:
                p.font.color.rgb = RGBColor(255, 165, 0)  # Orange
            else:
                p.font.color.rgb = RGBColor(255, 0, 0)  # Red
        
        # Add findings (limit to top 5 per slide)
        findings = analysis['findings'][:5]
        for finding in findings:
            p = text_frame.add_paragraph()
            severity = finding.get('severity', 'medium').upper()
            issue = finding.get('issue', 'N/A')
            p.text = f"[{severity}] {issue}"
            p.font.size = Pt(11)
            p.level = 1
            
            # Color code by severity
            if severity == 'CRITICAL':
                p.font.color.rgb = RGBColor(255, 0, 0)
            elif severity == 'HIGH':
                p.font.color.rgb = RGBColor(255, 140, 0)
            elif severity == 'MEDIUM':
                p.font.color.rgb = RGBColor(255, 215, 0)
        
        logger.info(f"Added findings slide for {resource_type}")

    def add_recommendations_slide(self, resource_type: str, analysis: Dict[str, Any]):
        """Add slide for recommendations of a specific resource type."""
        if not analysis or 'findings' not in analysis or not analysis['findings']:
            return
        
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{resource_type.replace('_', ' ').title()} - Recommendations"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add recommendations (limit to top 5 per slide)
        findings = analysis['findings'][:5]
        for i, finding in enumerate(findings):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
                
            resource = finding.get('resource', 'General')
            recommendation = finding.get('recommendation', 'N/A')
            p.text = f"{resource}: {recommendation}"
            p.font.size = Pt(11)
            p.level = 0
        
        logger.info(f"Added recommendations slide for {resource_type}")

    def add_best_practices_slide(self, resource_type: str, analysis: Dict[str, Any]):
        """Add slide for best practices met."""
        if not analysis or 'best_practices_met' not in analysis:
            return
        
        best_practices = analysis.get('best_practices_met', [])
        if not best_practices:
            return
        
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{resource_type.replace('_', ' ').title()} - Best Practices Met ✓"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, practice in enumerate(best_practices[:7]):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = f"✓ {practice}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0, 128, 0)
            p.level = 0
        
        logger.info(f"Added best practices slide for {resource_type}")

    def add_cost_analysis_slide(self, cost_analysis: Dict[str, Any]):
        """Add cost analysis and optimization recommendations slide."""
        if not cost_analysis:
            return
        
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Cost Analysis & Optimization"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        summary = cost_analysis.get('summary', {})
        opportunities = cost_analysis.get('optimization_opportunities', {})
        
        # Summary
        p = text_frame.paragraphs[0]
        total_findings = summary.get('total_findings', 0)
        p.text = f"Cost Optimization Opportunities: {total_findings}"
        p.font.size = Pt(16)
        p.font.bold = True
        if total_findings > 0:
            p.font.color.rgb = RGBColor(255, 140, 0)  # Orange
        else:
            p.font.color.rgb = RGBColor(0, 128, 0)  # Green
        
        # Opportunity breakdown
        p = text_frame.add_paragraph()
        p.text = f"  • Immediate Actions: {opportunities.get('immediate_actions', 0)}"
        p.font.size = Pt(12)
        p.level = 1
        
        p = text_frame.add_paragraph()
        p.text = f"  • Reviews Needed: {opportunities.get('review_needed', 0)}"
        p.font.size = Pt(12)
        p.level = 1
        
        p = text_frame.add_paragraph()
        p.text = f"  • Best Practice Suggestions: {opportunities.get('best_practices', 0)}"
        p.font.size = Pt(12)
        p.level = 1
        
        # Top recommendations
        recommendations = cost_analysis.get('recommendations', [])
        if recommendations:
            p = text_frame.add_paragraph()
            p.text = ""  # Empty line
            
            p = text_frame.add_paragraph()
            p.text = "Top Recommendations:"
            p.font.size = Pt(14)
            p.font.bold = True
            
            for i, rec in enumerate(recommendations[:3]):  # Top 3 recommendations
                p = text_frame.add_paragraph()
                priority = rec.get('priority', 99)
                summary_text = rec.get('summary', 'N/A')[:50]
                # Use descriptive priority text instead of raw number
                priority_text = "High Priority" if priority <= 2 else ("Medium Priority" if priority <= 4 else "Low Priority")
                p.text = f"  [{priority_text}] {summary_text}"
                p.font.size = Pt(11)
                p.level = 1
                
                if priority <= 2:
                    p.font.color.rgb = RGBColor(255, 0, 0)  # Red
                elif priority <= 4:
                    p.font.color.rgb = RGBColor(255, 140, 0)  # Orange
        
        logger.info("Added cost analysis slide")

    def add_cost_recommendations_slide(self, cost_analysis: Dict[str, Any]):
        """Add detailed cost recommendations slide."""
        recommendations = cost_analysis.get('recommendations', [])
        if not recommendations:
            return
        
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Cost Optimization Recommendations"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, rec in enumerate(recommendations[:5]):  # Top 5 recommendations
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            
            action = rec.get('action', 'N/A')[:70]
            affected = rec.get('affected_resources', 0)
            impact = rec.get('potential_impact', '')[:40]
            
            p.text = f"{rec.get('summary', 'N/A')[:60]}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.level = 0
            
            # Add action
            p = text_frame.add_paragraph()
            p.text = f"Action: {action}"
            p.font.size = Pt(10)
            p.level = 1
            
            # Add impact and affected count
            p = text_frame.add_paragraph()
            p.text = f"Affected: {affected} resources | Impact: {impact}"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.level = 1
        
        logger.info("Added cost recommendations slide")

    def add_tag_compliance_overview_slide(self, tag_analysis: Dict[str, Any]):
        """Add tag compliance overview slide."""
        if not tag_analysis:
            return
        
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Tag Compliance Overview"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        summary = tag_analysis.get('summary', {})
        compliance_rate = summary.get('overall_compliance_rate', 0)
        
        # Overall compliance
        p = text_frame.paragraphs[0]
        p.text = f"Overall Tag Compliance: {compliance_rate}%"
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Color code the compliance rate
        if compliance_rate >= 90:
            p.font.color.rgb = RGBColor(0, 128, 0)  # Green
        elif compliance_rate >= 70:
            p.font.color.rgb = RGBColor(255, 165, 0)  # Orange
        else:
            p.font.color.rgb = RGBColor(255, 0, 0)  # Red
        
        # Summary statistics
        p = text_frame.add_paragraph()
        p.text = ""  # Empty line
        
        stats = [
            f"Total Resources Analyzed: {summary.get('total_resources', 0)}",
            f"Resources with Tags: {summary.get('resources_with_tags', 0)}",
            f"Resources without Tags: {summary.get('resources_without_tags', 0)}",
            f"Required Tags: {summary.get('required_tags_count', 0)}",
            f"Resource Groups: {summary.get('total_resource_groups', 0)}"
        ]
        
        for stat in stats:
            p = text_frame.add_paragraph()
            p.text = stat
            p.font.size = Pt(12)
            p.level = 1
        
        # Required tags compliance summary
        required_tags = tag_analysis.get('required_tags_compliance', [])
        if required_tags:
            p = text_frame.add_paragraph()
            p.text = ""  # Empty line
            
            p = text_frame.add_paragraph()
            p.text = "Required Tags Compliance:"
            p.font.size = Pt(14)
            p.font.bold = True
            
            for tag_info in required_tags[:5]:  # Top 5 required tags
                tag_name = tag_info.get('tag_name', 'Unknown')
                tag_compliance = tag_info.get('compliance_percentage', 0)
                
                p = text_frame.add_paragraph()
                p.text = f"  {tag_name}: {tag_compliance}%"
                p.font.size = Pt(11)
                p.level = 1
                
                if tag_compliance >= 90:
                    p.font.color.rgb = RGBColor(0, 128, 0)  # Green
                elif tag_compliance >= 70:
                    p.font.color.rgb = RGBColor(255, 165, 0)  # Orange
                else:
                    p.font.color.rgb = RGBColor(255, 0, 0)  # Red
        
        logger.info("Added tag compliance overview slide")

    def add_tag_compliance_by_resource_group_slide(self, tag_analysis: Dict[str, Any]):
        """Add tag compliance by resource group slide."""
        resource_groups_details = tag_analysis.get('resource_groups_details', [])
        if not resource_groups_details:
            return
        
        # Create multiple slides if needed (max 5 resource groups per slide)
        for slide_num, i in enumerate(range(0, len(resource_groups_details), 5)):
            slide_layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            if len(resource_groups_details) > 5:
                title.text = f"Tag Compliance by Resource Group ({slide_num + 1})"
            else:
                title.text = "Tag Compliance by Resource Group"
            
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            
            rg_slice = resource_groups_details[i:i+5]
            first = True
            
            for rg in rg_slice:
                if not first:
                    p = text_frame.add_paragraph()
                    p.text = ""  # Empty line
                else:
                    first = False
                
                rg_name = rg.get('name', 'Unknown')
                rg_compliance = rg.get('compliance_rate', 0)
                rg_missing_tags = rg.get('missing_tags', [])
                non_compliant_resources = rg.get('non_compliant_resources', 0)
                total_resources = rg.get('total_resources', 0)
                
                # Resource group name and compliance
                p = text_frame.add_paragraph() if not first else text_frame.paragraphs[0]
                p.text = f"{rg_name} - {rg_compliance}% compliant"
                p.font.size = Pt(12)
                p.font.bold = True
                p.level = 0
                
                if rg_compliance >= 90:
                    p.font.color.rgb = RGBColor(0, 128, 0)
                elif rg_compliance >= 70:
                    p.font.color.rgb = RGBColor(255, 165, 0)
                else:
                    p.font.color.rgb = RGBColor(255, 0, 0)
                
                # Resource group missing tags
                if rg_missing_tags:
                    p = text_frame.add_paragraph()
                    p.text = f"  RG Missing Tags: {', '.join(rg_missing_tags)}"
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(255, 0, 0)
                    p.level = 1
                else:
                    p = text_frame.add_paragraph()
                    p.text = "  RG Tags: ✓ All required tags present"
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(0, 128, 0)
                    p.level = 1
                
                # Resource compliance summary
                p = text_frame.add_paragraph()
                p.text = f"  Resources: {non_compliant_resources} of {total_resources} missing tags"
                p.font.size = Pt(10)
                p.level = 1
                if non_compliant_resources > 0:
                    p.font.color.rgb = RGBColor(255, 140, 0)
        
        logger.info("Added tag compliance by resource group slides")

    def add_summary_slide(self):
        """Add final summary slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Next Steps"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        next_steps = [
            "Review all critical and high-severity findings",
            "Prioritize security-related issues",
            "Review the improvement backlog",
            "Implement recommended changes systematically",
            "Schedule regular assessments",
            "Monitor ongoing compliance"
        ]
        
        for i, step in enumerate(next_steps):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = f"{i+1}. {step}"
            p.font.size = Pt(14)
            p.level = 0
        
        logger.info("Added summary slide")

    def generate_report(self, resources: Dict[str, List[Dict[str, Any]]], 
                       analyses: Dict[str, Any], output_path: str):
        """Generate complete PowerPoint report."""
        logger.info("Generating PowerPoint report...")
        
        # Add slides
        self.add_title_slide()
        
        # Executive summary
        if 'executive_summary' in analyses:
            self.add_executive_summary_slide(analyses['executive_summary'])
        
        # Resource overview
        self.add_resource_overview_slide(resources)
        
        # Cost analysis (if available)
        if 'cost_analysis' in analyses:
            self.add_cost_analysis_slide(analyses['cost_analysis'])
            self.add_cost_recommendations_slide(analyses['cost_analysis'])
        
        # Tag analysis (if available)
        if 'tag_analysis' in analyses:
            self.add_tag_compliance_overview_slide(analyses['tag_analysis'])
            self.add_tag_compliance_by_resource_group_slide(analyses['tag_analysis'])
        
        # Add slides for each resource type
        resource_types = [
            'virtual_machines',
            'storage_accounts', 
            'network_security_groups',
            'virtual_networks'
        ]
        
        for resource_type in resource_types:
            if resource_type in analyses:
                analysis = analyses[resource_type]
                self.add_findings_slide(resource_type, analysis)
                self.add_recommendations_slide(resource_type, analysis)
                self.add_best_practices_slide(resource_type, analysis)
        
        # Add slides for generic resources grouped by type
        if 'generic_resources' in analyses:
            generic_analysis = analyses['generic_resources']
            if 'resource_types' in generic_analysis:
                for resource_type, analysis in generic_analysis['resource_types'].items():
                    # Use simplified resource type name for slide title
                    simple_type = resource_type.split('/')[-1] if '/' in resource_type else resource_type
                    self.add_findings_slide(simple_type, analysis)
                    self.add_recommendations_slide(simple_type, analysis)
                    self.add_best_practices_slide(simple_type, analysis)
        
        # Final summary
        self.add_summary_slide()
        
        # Save presentation
        self.prs.save(output_path)
        logger.info(f"PowerPoint report saved to: {output_path}")
